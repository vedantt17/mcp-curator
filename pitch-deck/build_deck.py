"""
MCP Curator pitch deck — 6 slides, 16:9, Zinc + Emerald palette.
Run: python build_deck.py  (writes ./mcp-curator-pitch.pptx)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 16:9 in inches: 13.333 x 7.5
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0xFA, 0xFA, 0xFA)
INK = RGBColor(0x18, 0x18, 0x1B)
MUTED = RGBColor(0x52, 0x52, 0x5B)
SUBTLE = RGBColor(0xA1, 0xA1, 0xAA)
BORDER = RGBColor(0xE4, 0xE4, 0xE7)
ACCENT = RGBColor(0x05, 0x96, 0x69)
ACCENT_SOFT = RGBColor(0xA7, 0xF3, 0xD0)
SURFACE = RGBColor(0xF4, 0xF4, 0xF5)


def blank_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def add_text(slide, text, left, top, width, height, *, size=16, bold=False, color=INK,
              font="Helvetica", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_serif=False):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Georgia" if font_serif else font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tx


def add_runs(slide, runs, left, top, width, height, *, size=16, font="Helvetica",
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_serif=False, line_spacing=None):
    """runs: list of (text, {bold, color, size}) tuples. None text means new paragraph."""
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    para = tf.paragraphs[0]
    para.alignment = align
    if line_spacing:
        para.line_spacing = line_spacing
    for item in runs:
        if item is None:
            para = tf.add_paragraph()
            para.alignment = align
            if line_spacing:
                para.line_spacing = line_spacing
            continue
        text, opts = item
        r = para.add_run()
        r.text = text
        r.font.name = opts.get("font", "Georgia" if font_serif else font)
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", False)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", INK)
    return tx


def add_rule(slide, x, y, w, color=BORDER, thickness=0.75):
    """Thin horizontal rule."""
    h = Pt(thickness)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False
    return line


def add_box(slide, x, y, w, h, *, fill=None, line=None, line_w=0.75, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(line_w)
    box.shadow.inherit = False
    if radius:
        # 0..1 of half min dimension; small corner
        box.adjustments[0] = 0.08
    return box


def page_chrome(slide, page_label, page_num):
    """Header logo + page number footer."""
    add_runs(slide, [
        ("mcp", {"size": 14, "bold": True, "color": INK}),
        ("-", {"size": 14, "bold": True, "color": ACCENT}),
        ("curator", {"size": 14, "bold": True, "color": INK}),
    ], Inches(0.5), Inches(0.3), Inches(4), Inches(0.4), font_serif=True)

    add_text(slide, page_label, Inches(8.833), Inches(0.3), Inches(4), Inches(0.4),
             size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(slide, f"{page_num} / 6", Inches(0.5), Inches(7.05), Inches(12.333), Inches(0.3),
             size=9, color=SUBTLE, align=PP_ALIGN.RIGHT)


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ────────────────────────────────────────────────────────────────────────────
s = blank_slide()

# Big logotype
add_runs(s, [
    ("mcp", {"size": 88, "bold": True, "color": INK}),
    ("-", {"size": 88, "bold": True, "color": ACCENT}),
    ("curator", {"size": 88, "bold": True, "color": INK}),
], Inches(0.5), Inches(2.0), Inches(12.333), Inches(1.5), font_serif=True)

# Headline
add_text(s,
         "Turn any OpenAPI spec into a Claude-ready MCP server.",
         Inches(0.5), Inches(3.7), Inches(12.333), Inches(1.0),
         size=30, color=INK)

# Subtitle
add_runs(s, [
    ("Paste a spec. ", {"size": 18, "color": MUTED}),
    ("Claude curates the endpoints.", {"size": 18, "color": INK}),
    (" Install with ", {"size": 18, "color": MUTED}),
    ("npx", {"size": 18, "color": INK, "font": "Courier New"}),
    (". 60 seconds.", {"size": 18, "color": MUTED}),
], Inches(0.5), Inches(4.55), Inches(12.333), Inches(0.5))

# Bottom rule + meta
add_rule(s, Inches(0.5), Inches(6.6), Inches(12.333))
add_text(s, "SF · Claude for Technical Founders · 2026",
         Inches(0.5), Inches(6.75), Inches(8), Inches(0.4),
         size=11, color=MUTED)
add_text(s, "Tanmay  ·  Vedant",
         Inches(8.5), Inches(6.75), Inches(4.333), Inches(0.4),
         size=11, color=MUTED, align=PP_ALIGN.RIGHT)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Problem
# ────────────────────────────────────────────────────────────────────────────
s = blank_slide()
page_chrome(s, "the problem", 2)

add_text(s, "Mechanical OpenAPI → MCP converters",
         Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.7),
         size=36, color=INK, font_serif=True)
add_text(s, "give you 200 dumb tools.",
         Inches(0.5), Inches(1.95), Inches(12.333), Inches(0.7),
         size=36, color=ACCENT, font_serif=True)

# Left column: simulated mechanical dump
col_x = Inches(0.5)
col_y = Inches(3.4)
col_w = Inches(6.0)
col_h = Inches(3.0)
add_box(s, col_x, col_y, col_w, col_h, fill=SURFACE, line=BORDER)

mechanical_label = add_text(s, "openapi-mcp-server  /  Speakeasy  /  Stainless",
                             Inches(0.7), Inches(3.55), Inches(5.7), Inches(0.3),
                             size=10, color=MUTED, font="Courier New")

# Show many tool names faintly
mechanical_names = [
    "post_emails  patch_emails_by_id  delete_emails_by_id",
    "post_emails_batch  get_emails_by_id  get_emails_received",
    "post_audiences  get_audiences  delete_audiences_by_id",
    "post_audiences_id_contacts  get_audiences_id_contacts",
    "delete_audiences_id_contacts_by_id  patch_contacts_by_id",
    "post_contacts  get_contacts  delete_contacts_by_id",
    "post_domains  get_domains  delete_domains_by_id",
    "post_domains_by_id_verify  get_domain_records",
    "post_broadcasts  get_broadcasts  get_broadcasts_by_id",
    "post_broadcasts_by_id_send  delete_broadcasts_by_id",
    "post_api_keys  get_api_keys  delete_api_keys_by_id",
    "… (47 more endpoints, each as its own tool)",
]
y = 4.0
for line in mechanical_names:
    is_last = "…" in line
    add_text(s, line, Inches(0.7), Inches(y), col_w - Inches(0.4), Inches(0.22),
             size=9, color=SUBTLE if is_last else MUTED, font="Courier New")
    y += 0.2

# Right column: the consequences
right_x = Inches(7.0)
right_w = Inches(5.833)

add_text(s, "The cost", Inches(7.0), Inches(3.4), right_w, Inches(0.4),
         size=14, bold=True, color=INK)

bullets = [
    ("Naming overlap.", " update_pet vs update_pet_with_form — same intent, different shape."),
    ("Cryptic descriptions.", " OpenAPI specs are written for SDKs, not LLMs."),
    ("Destructive endpoints exposed.", " The model can DELETE prod data on a hallucination."),
    ("Pagination, headers, $refs.", " Implementation details leak into the tool surface."),
]
y = 3.95
for head, tail in bullets:
    add_runs(s, [
        ("→  ", {"size": 12, "color": ACCENT, "bold": True}),
        (head, {"size": 12, "color": INK, "bold": True}),
        (tail, {"size": 12, "color": MUTED}),
    ], Inches(7.0), Inches(y), right_w, Inches(0.5), line_spacing=1.25)
    y += 0.55

# Bottom takeaway
add_rule(s, Inches(0.5), Inches(6.5), Inches(12.333), color=BORDER)
add_text(s, "Result: the LLM is overwhelmed. Tool calls degrade. The integration looks clever but performs worse than no integration at all.",
         Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.4),
         size=11, color=MUTED, align=PP_ALIGN.LEFT)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — The fix
# ────────────────────────────────────────────────────────────────────────────
s = blank_slide()
page_chrome(s, "the fix", 3)

add_text(s, "We let Claude curate.",
         Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.8),
         size=42, color=INK, font_serif=True)
add_text(s, "≤12 LLM-ready tools. Claude picks them, names them, and writes their descriptions.",
         Inches(0.5), Inches(2.05), Inches(12.333), Inches(0.6),
         size=16, color=MUTED)

# Four cards
cards = [
    ("Picks the useful ones.",
     "Out of 200 endpoints, the ≤12 most LLM-actionable. Read-over-write when both qualify."),
    ("Drops what's dangerous.",
     "DELETEs and irreversible POSTs are filtered by default. The model can't wipe prod by accident."),
    ("Merges related ops.",
     "GET /pets/by-status + GET /pets/by-tag → one find_pets tool with a flexible input schema."),
    ("Rewrites descriptions for LLMs.",
     "PURPOSE / WHEN to use / WHEN NOT to use. Plus an example invocation per tool."),
]

card_y = 2.95
card_h = Inches(1.8)
gap = Inches(0.2)
total_w = Inches(12.333)
card_w = Emu(int((total_w - gap * 3) / 4))

for i, (title, body) in enumerate(cards):
    x = Inches(0.5) + (card_w + gap) * i
    add_box(s, x, Inches(card_y), card_w, card_h, fill=SURFACE, line=BORDER, radius=True)
    add_text(s, title, x + Inches(0.25), Inches(card_y + 0.2), card_w - Inches(0.5), Inches(0.6),
             size=14, bold=True, color=INK)
    add_text(s, body, x + Inches(0.25), Inches(card_y + 0.7), card_w - Inches(0.5), Inches(1.0),
             size=11, color=MUTED)

# Tagline
add_text(s, "Curate, don't translate.",
         Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.6),
         size=20, color=ACCENT, font_serif=True)

# Boundary explanation
add_runs(s, [
    ("The split is deliberate. ", {"size": 12, "color": INK, "bold": True}),
    ("Claude does the judgment work — choosing endpoints, writing descriptions, merging schemas. Codegen does the deterministic work — emitting fetch handlers from the chosen ops. ", {"size": 12, "color": MUTED}),
    ("Hallucination risk is contained to descriptions; the network calls are correct by construction.", {"size": 12, "color": INK}),
], Inches(0.5), Inches(6.0), Inches(12.333), Inches(1.0), line_spacing=1.3)

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — How it works
# ────────────────────────────────────────────────────────────────────────────
s = blank_slide()
page_chrome(s, "how it works", 4)

add_text(s, "How it works",
         Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.7),
         size=36, color=INK, font_serif=True)
add_text(s, "Two phases: generation (server-side) and install (user's machine).",
         Inches(0.5), Inches(1.95), Inches(12.333), Inches(0.4),
         size=14, color=MUTED)

# Phase 1: server-side
phase1_y = Inches(2.7)
add_text(s, "1.  Generation",
         Inches(0.5), phase1_y, Inches(4), Inches(0.4),
         size=13, bold=True, color=INK)
add_text(s, "(server, ~30s with prompt caching)",
         Inches(2.4), phase1_y + Inches(0.04), Inches(6), Inches(0.4),
         size=11, color=MUTED)

box_y = phase1_y + Inches(0.5)
box_h = Inches(0.65)
box_w = Inches(2.7)
gap_x = Inches(0.2)
arrow_w = Inches(0.25)

steps_p1 = [
    ("Paste OpenAPI", "JSON or YAML, URL or inline"),
    ("Claude curates", "≤12 tools, descriptions rewritten"),
    ("Codegen", "deterministic fetch handlers"),
    ("Store + return id", "Upstash, returns nanoid"),
]
x = Inches(0.5)
for i, (t, sub) in enumerate(steps_p1):
    add_box(s, x, box_y, box_w, box_h, fill=SURFACE, line=BORDER, radius=True)
    add_text(s, t, x + Inches(0.15), box_y + Inches(0.08), box_w - Inches(0.3), Inches(0.3),
             size=12, bold=True, color=INK)
    add_text(s, sub, x + Inches(0.15), box_y + Inches(0.36), box_w - Inches(0.3), Inches(0.3),
             size=9, color=MUTED)
    if i < 3:
        ax = x + box_w + Inches(0.02)
        add_text(s, "→", ax, box_y + Inches(0.18), arrow_w, Inches(0.4),
                 size=18, color=ACCENT, align=PP_ALIGN.CENTER)
    x = x + box_w + arrow_w + Emu(0)

# Phase 2: client-side
phase2_y = Inches(4.4)
add_text(s, "2.  Install",
         Inches(0.5), phase2_y, Inches(4), Inches(0.4),
         size=13, bold=True, color=INK)
add_text(s, "(user's machine, single command)",
         Inches(1.85), phase2_y + Inches(0.04), Inches(6), Inches(0.4),
         size=11, color=MUTED)

# command + flow
cmd_box_y = phase2_y + Inches(0.5)
cmd_box_h = Inches(0.65)

steps_p2 = [
    ("npx mcp-from-spec id", "fetch generated code"),
    ("npm install + tsx", "spawns over stdio"),
    ("Claude Desktop", "tools appear, ready to call"),
]
box_w2 = Inches(3.7)
x = Inches(0.5)
for i, (t, sub) in enumerate(steps_p2):
    add_box(s, x, cmd_box_y, box_w2, cmd_box_h, fill=SURFACE, line=BORDER, radius=True)
    add_text(s, t, x + Inches(0.15), cmd_box_y + Inches(0.08), box_w2 - Inches(0.3), Inches(0.3),
             size=12, bold=True, color=INK, font="Courier New" if i == 0 else "Helvetica")
    add_text(s, sub, x + Inches(0.15), cmd_box_y + Inches(0.36), box_w2 - Inches(0.3), Inches(0.3),
             size=9, color=MUTED)
    if i < 2:
        ax = x + box_w2 + Inches(0.05)
        add_text(s, "→", ax, cmd_box_y + Inches(0.18), arrow_w, Inches(0.4),
                 size=18, color=ACCENT, align=PP_ALIGN.CENTER)
    x = x + box_w2 + arrow_w + Inches(0.05)

# Stack note
add_rule(s, Inches(0.5), Inches(6.0), Inches(12.333))
add_runs(s, [
    ("Stack:  ", {"size": 11, "bold": True, "color": INK}),
    ("Next.js 16  ·  Anthropic SDK + prompt caching  ·  Upstash Redis  ·  ", {"size": 11, "color": MUTED}),
    ("@modelcontextprotocol/sdk", {"size": 11, "color": INK, "font": "Courier New"}),
    (" (stdio)  ·  Vercel", {"size": 11, "color": MUTED}),
], Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.5))

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Numbers
# ────────────────────────────────────────────────────────────────────────────
s = blank_slide()
page_chrome(s, "numbers", 5)

add_text(s, "Three real specs, no prompt tweaks.",
         Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.7),
         size=36, color=INK, font_serif=True)
add_text(s, "Each row: input ops on the left, what Claude actually returned on the right.",
         Inches(0.5), Inches(1.95), Inches(12.333), Inches(0.4),
         size=14, color=MUTED)

# Table-like rows
rows = [
    ("Petstore",          "Swagger demo",                  "19", "12", "list_pets, find_pets, place_order, ..."),
    ("Scalar API",        "OpenAPI 3.1, real production",  "41", "12", "list_apis, get_api_document, manage_access_group (3 ops merged), ..."),
    ("Resend",            "Email, 83 endpoints",           "83", "12", "send_email, send_batch_emails, manage_contact (3 ops merged), ..."),
]

table_y = 2.85
row_h = Inches(1.05)
row_gap = Inches(0.15)

for i, (name, sub, input_n, output_n, examples) in enumerate(rows):
    y = table_y + (i * 1.2)
    add_box(s, Inches(0.5), Inches(y), Inches(12.333), row_h, fill=SURFACE, line=BORDER, radius=True)

    # Spec name
    add_text(s, name, Inches(0.8), Inches(y + 0.18), Inches(2.5), Inches(0.4),
             size=18, bold=True, color=INK, font_serif=True)
    add_text(s, sub, Inches(0.8), Inches(y + 0.62), Inches(2.5), Inches(0.4),
             size=10, color=MUTED)

    # Input number
    add_text(s, input_n, Inches(3.5), Inches(y + 0.15), Inches(1.0), Inches(0.7),
             size=36, bold=True, color=MUTED, font_serif=True, align=PP_ALIGN.CENTER)
    add_text(s, "raw ops", Inches(3.5), Inches(y + 0.7), Inches(1.0), Inches(0.3),
             size=9, color=SUBTLE, align=PP_ALIGN.CENTER)

    # Arrow
    add_text(s, "→", Inches(4.6), Inches(y + 0.25), Inches(0.5), Inches(0.5),
             size=24, color=ACCENT, align=PP_ALIGN.CENTER)

    # Output number
    add_text(s, output_n, Inches(5.1), Inches(y + 0.15), Inches(1.0), Inches(0.7),
             size=36, bold=True, color=ACCENT, font_serif=True, align=PP_ALIGN.CENTER)
    add_text(s, "curated", Inches(5.1), Inches(y + 0.7), Inches(1.0), Inches(0.3),
             size=9, color=SUBTLE, align=PP_ALIGN.CENTER)

    # Sample tools
    add_text(s, "tools", Inches(6.4), Inches(y + 0.18), Inches(0.6), Inches(0.3),
             size=10, color=SUBTLE)
    add_text(s, examples, Inches(6.4), Inches(y + 0.45), Inches(5.7), Inches(0.6),
             size=10, color=MUTED, font="Courier New")

# ────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Demo / closing
# ────────────────────────────────────────────────────────────────────────────
s = blank_slide()
page_chrome(s, "demo", 6)

add_text(s, "See it.",
         Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.8),
         size=56, color=INK, font_serif=True)

add_text(s, "Paste an OpenAPI URL. We'll do the rest live.",
         Inches(0.5), Inches(2.1), Inches(12.333), Inches(0.5),
         size=18, color=MUTED)

# Two boxes: the URL + the install command
add_box(s, Inches(0.5), Inches(3.0), Inches(6.0), Inches(2.5), fill=SURFACE, line=BORDER, radius=True)
add_text(s, "the demo URL", Inches(0.75), Inches(3.2), Inches(5.5), Inches(0.3),
         size=10, bold=True, color=ACCENT)
add_text(s, "mcp-curator.vercel.app", Inches(0.75), Inches(3.55), Inches(5.5), Inches(0.6),
         size=22, bold=True, color=INK, font="Courier New")
add_text(s, "Paste any spec.  Get a Claude-curated MCP server.  Install in one command.",
         Inches(0.75), Inches(4.4), Inches(5.5), Inches(1.0),
         size=11, color=MUTED)

add_box(s, Inches(6.833), Inches(3.0), Inches(6.0), Inches(2.5), fill=SURFACE, line=BORDER, radius=True)
add_text(s, "the install", Inches(7.083), Inches(3.2), Inches(5.5), Inches(0.3),
         size=10, bold=True, color=ACCENT)
add_text(s, "npx mcp-from-spec <id>", Inches(7.083), Inches(3.55), Inches(5.5), Inches(0.6),
         size=20, bold=True, color=INK, font="Courier New")
add_text(s, "Drop into Claude Desktop config. Restart. Tools appear. The model can call your API.",
         Inches(7.083), Inches(4.4), Inches(5.5), Inches(1.0),
         size=11, color=MUTED)

# Closing line
add_text(s, "Every API on the web — Claude-callable in 60 seconds.",
         Inches(0.5), Inches(5.9), Inches(12.333), Inches(0.6),
         size=22, color=ACCENT, font_serif=True, align=PP_ALIGN.CENTER)

# Footer
add_rule(s, Inches(0.5), Inches(6.7), Inches(12.333))
add_runs(s, [
    ("github.com/vedantt17/mcp-curator", {"size": 11, "color": INK, "font": "Courier New"}),
    ("        ", {"size": 11, "color": MUTED}),
    ("Tanmay  ·  Vedant", {"size": 11, "color": MUTED}),
], Inches(0.5), Inches(6.85), Inches(12.333), Inches(0.4), align=PP_ALIGN.CENTER)


out = "mcp-curator-pitch.pptx"
prs.save(out)
print(f"wrote {out}")
